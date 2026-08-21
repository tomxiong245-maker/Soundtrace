#!/usr/bin/env python3
"""EP05 pre 讲稿 · 生成 .pptx (22 页 · 20 分钟 · 中文)

- 输入源: 8:21 讲稿 (主) + 8:20 流程图源 (只挖流程图)
- 输出: EP05-pre-2026-8-21.pptx (同目录)
- 严格按 presentation-production SKILL Phase 5 排版 + Phase 6 speaker notes
- 直接可念 · 无 meta 开头 · 每边界 transition 只一次

Run:
  python3 build_pptx.py
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

OUT = Path(__file__).parent / "EP05-pre-2026-8-21.pptx"

# ================== 色板 · 匹配 html 版一致 ==================
C_BG_LIGHT = RGBColor(0xFA, 0xFA, 0xF9)  # 温白
C_BG_DARK  = RGBColor(0x0F, 0x17, 0x2A)  # 深蓝黑
C_INK      = RGBColor(0x0F, 0x17, 0x2A)  # 主文本
C_INK2     = RGBColor(0x29, 0x25, 0x24)  # 次
C_INK3     = RGBColor(0x57, 0x53, 0x4E)  # 三
C_GRAY     = RGBColor(0xA8, 0xA2, 0x9E)  # HUD
C_ACCENT   = RGBColor(0xDC, 0x26, 0x26)  # 红 · accent
C_GREEN    = RGBColor(0x05, 0x96, 0x69)  # 通过
C_BLUE     = RGBColor(0x1E, 0x40, 0xAF)  # 准备层
C_PURPLE   = RGBColor(0x7C, 0x3A, 0xED)  # PREFERENCE
C_ORANGE   = RGBColor(0xC2, 0x41, 0x0C)  # PARAMETER
C_DGREEN   = RGBColor(0x16, 0x65, 0x34)  # QA
C_CARD     = RGBColor(0xFF, 0xFF, 0xFF)
C_YELLOW_BG = RGBColor(0xFE, 0xF3, 0xC7)  # main-decider 底
C_YELLOW_BORDER = RGBColor(0xF5, 0x9E, 0x0B)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xE7, 0xE5, 0xE4)

FONT = "PingFang SC"
FONT_MONO = "Menlo"

# ================== 幻灯片尺寸 · 16:9 (Widescreen 13.333 x 7.5 in) ==================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]  # blank layout

# ================== 工具函数 ==================
def add_slide(bg=C_BG_LIGHT):
    s = prs.slides.add_slide(BLANK)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    return s

def add_text(slide, left, top, width, height, text, *,
             size=20, bold=False, color=C_INK, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.35):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = str(text).split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_rect(slide, left, top, width, height, *, fill=C_CARD, line=None, line_w=0):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line:
        r.line.color.rgb = line
        r.line.width = Pt(line_w or 1)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r

def add_left_border(slide, left, top, w_border, height, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w_border, height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r

def add_top_border(slide, left, top, width, h_border, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h_border)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r

def add_sec_chip(slide, text, top=Inches(0.55)):
    """§ N · xxx chip · 左上"""
    chip = add_rect(slide, Inches(0.75), top, Inches(2.6), Inches(0.35), fill=RGBColor(0xFE, 0xF2, 0xF2))
    add_text(slide, Inches(0.75), top, Inches(2.6), Inches(0.35), text,
             size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def add_h2(slide, text, *, top=Inches(1.0), color=C_INK):
    add_text(slide, Inches(0.75), top, Inches(11.83), Inches(0.9), text,
             size=32, bold=True, color=color)

_DARK_PAGES = set()  # 1-indexed · 记录哪些页要暗色 HUD (封面 / 收尾)

def add_page_num(slide, n=None, total=None, dark=False):
    """Deferred · 数字忽略 · finalize_page_nums 统一写"""
    if dark:
        _DARK_PAGES.add(len(prs.slides))  # 当前 slide 是 1-index 最后一张

def finalize_page_nums():
    """所有 slide 加完后 · 统一写页码 (n / total 用实际总数)"""
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        is_dark = i in _DARK_PAGES
        color = RGBColor(0x94, 0xA3, 0xB8) if is_dark else C_GRAY
        add_text(slide, Inches(0.75), Inches(7.1), Inches(2), Inches(0.3),
                 f"{i} / {total}", size=10, color=color)
        add_text(slide, Inches(9.5), Inches(7.1), Inches(3.5), Inches(0.3),
                 "Soundtrace · 2026-08-21 · 熊镇正", size=10, color=color, align=PP_ALIGN.RIGHT)

def add_audio(slide, mp3_path, left, top, width=Inches(0.55), height=Inches(0.55)):
    """嵌入音频 · Keynote/PowerPoint 播放时点扬声器图标.
    若 mp3 找不到 (clone 后无媒体) · 优雅降级为占位圆点 + 标注."""
    from pathlib import Path as _P
    if not _P(str(mp3_path)).exists():
        print(f"[audio missing · placeholder] {mp3_path}", flush=True)
        r = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        r.fill.solid(); r.fill.fore_color.rgb = C_GRAY
        r.line.fill.background(); r.shadow.inherit = False
        add_text(slide, left, top, width, height, "♪", size=18, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return None
    try:
        return slide.shapes.add_movie(
            str(mp3_path), left, top, width, height,
            poster_frame_image=None, mime_type='audio/mpeg'
        )
    except Exception as e:
        r = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        r.fill.solid(); r.fill.fore_color.rgb = C_ACCENT
        r.line.fill.background()
        add_text(slide, left, top, width, height, "▶", size=18, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        print(f"[audio embed fail] {mp3_path} · {e}", flush=True)
        return None

def add_picture_safe(slide, img_path, left, top, width=None, height=None):
    """插图 · 若文件找不到 · 用带边框的占位框 + 文件名."""
    from pathlib import Path as _P
    if not _P(str(img_path)).exists():
        w = width or Inches(3); h = height or Inches(2)
        add_rect(slide, left, top, w, h, fill=RGBColor(0xF5, 0xF5, 0xF4),
                 line=C_GRAY, line_w=1)
        add_text(slide, left, top, w, h,
                 f"[缺失图片 · 请补]\n{_P(str(img_path)).name}",
                 size=14, color=C_INK3, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        print(f"[image missing · placeholder] {img_path}", flush=True)
        return None
    return slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)

def add_page_num_num(slide, n, total=22, dark=False):
    color = RGBColor(0x94, 0xA3, 0xB8) if dark else C_GRAY
    add_text(slide, Inches(0.75), Inches(7.1), Inches(2), Inches(0.3),
             f"{n} / {total}", size=10, color=color)
    add_text(slide, Inches(9.5), Inches(7.1), Inches(3.5), Inches(0.3),
             "Soundtrace · 2026-08-21 · 熊镇正", size=10, color=color, align=PP_ALIGN.RIGHT)

def add_notes(slide, text):
    """Speaker notes · 打印视图 + Presenter View 可见"""
    notes = slide.notes_slide.notes_text_frame
    notes.text = ""
    for i, line in enumerate(text.strip().split("\n")):
        p = notes.paragraphs[0] if i == 0 else notes.add_paragraph()
        r = p.add_run()
        r.text = line.strip()
        r.font.name = FONT
        r.font.size = Pt(14)

# ================== 流程图 stage card ==================
def add_flow_stage(slide, top, stage_id, body_html, *, color=C_BLUE, main=False, gate=False):
    """一个 stage card · 返回下一 top"""
    left = Inches(1.1)
    width = Inches(11.13)
    height = Inches(0.68) if not main else Inches(0.85)
    fill = C_YELLOW_BG if main else (RGBColor(0xFA, 0xFA, 0xF9) if gate else C_CARD)
    # 卡片背景
    add_rect(slide, left, top, width, height, fill=fill,
             line=C_YELLOW_BORDER if main else None, line_w=2 if main else 0)
    # 左 border
    add_left_border(slide, left, top, Inches(0.06), height, color)
    # stage-id
    add_text(slide, left + Inches(0.15), top, Inches(1.7), height, stage_id,
             size=11, bold=True, color=C_INK3, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
    # body
    add_text(slide, left + Inches(1.9), top, width - Inches(2.1), height, body_html,
             size=14, color=C_INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    return top + height + Inches(0.08)

def add_arrow_down(slide, top, note=""):
    """向下箭头"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.55), top, Inches(0.25), Inches(0.22))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = C_GRAY
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    if note:
        add_text(slide, Inches(7.0), top, Inches(4), Inches(0.22), note,
                 size=10, color=C_INK3, anchor=MSO_ANCHOR.MIDDLE)
    return top + Inches(0.28)

def add_brand_line(slide, left, top, width=Inches(1.5), height=Inches(0.05)):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid(); r.fill.fore_color.rgb = C_ACCENT
    r.line.fill.background()
    r.shadow.inherit = False

# =========================================================
# ==================== 22 页 slides =======================
# =========================================================

# ---- P01 · 封面 ----
s = add_slide(bg=C_BG_DARK)
add_brand_line(s, Inches(0.85), Inches(1.9))
add_text(s, Inches(0.85), Inches(2.1), Inches(11), Inches(1.2),
         "Soundtrace", size=64, bold=True, color=C_WHITE)
add_text(s, Inches(0.85), Inches(3.5), Inches(11), Inches(1.5),
         "章鱼 AI 播客个性化剪辑助手\n本地跑 · 每一步可追溯 · 真人拍板",
         size=22, color=RGBColor(0x94, 0xA3, 0xB8), line_spacing=1.5)
add_text(s, Inches(0.85), Inches(5.5), Inches(11), Inches(0.5),
         "熊镇正 · 香港中文大学（深圳）",
         size=18, color=RGBColor(0xCB, 0xD5, 0xE1))
add_page_num(s, 1, dark=True)
add_notes(s, """大家好 · 我是熊镇正 · 来自香港中文大学（深圳）。今天讲我一个人做的一套东西 —— Soundtrace · 章鱼 AI 播客个性化剪辑助手。它不是"一键出片"的黑盒子 · 是一套在自己电脑上跑、每一步都能看到在干什么、最后由真人拍板的多轨播客后期助手。它不追求你按一下就出成品 · 它追求把重复劳动接管掉、学会剪辑师的偏好 · 让每一次决定都变成下一期节目能用上的经验。我今天分五段讲。""")

# ---- P02 · 议程 ----
s = add_slide()
add_sec_chip(s, "议程 · 20 分钟")
add_h2(s, "今天分五段讲")
agenda = [
    ("01", "Demo · 听感对比", "A/B 两段音频"),
    ("02", "背景 · 为什么这件事有价值", "三条真正的难点"),
    ("03", '框架 · 从"能力封装"开始', "6 层 skill · 4 张流程图"),
    ("04", "踩坑经验 · 一周实录", "Optuna · LLM · 反思"),
    ("05", "未来 · 还能往哪里长", "偏好积累的四个方向"),
]
y = Inches(2.1)
for num, title, hint in agenda:
    add_rect(s, Inches(0.75), y, Inches(11.83), Inches(0.75), fill=C_CARD)
    add_left_border(s, Inches(0.75), y, Inches(0.06), Inches(0.75), C_ACCENT)
    add_text(s, Inches(1.0), y, Inches(1.2), Inches(0.75), num,
             size=22, bold=True, color=C_ACCENT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.2), y, Inches(7.5), Inches(0.75), title,
             size=18, bold=True, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.8), y, Inches(2.75), Inches(0.75), hint,
             size=12, color=C_INK3, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    y += Inches(0.85)
add_page_num(s, 2)
add_notes(s, """先放一段听感对比 · 再讲这件事为什么难。然后从整体流程切入我的项目框架 —— 我把它拆成 6 层 skill、4 张流程图。之后是一周的踩坑实录 —— 参数、候选、mentor 采访引发的反思都在这一段。最后聊一下未来还能往哪里长。整场大约 20 分钟。先听第一段。""")

# ---- P03 · Demo · 4 音频嵌入 ----
s = add_slide()
add_sec_chip(s, "§ 1 · Demo")
add_h2(s, "先放两段音频 · 请对比")
# 2 × 2 音频卡片
audios = [
    ("EXP019 · 原始",          "/Users/renting/Desktop/EXP019_original.mp3",       False),
    ("EXP019 · Soundtrace 剪", "/Users/renting/Desktop/EXP019_cut.mp3",            True),
    ("EXP002 · 原始",          "/Users/renting/Desktop/EXP002_original_wider.mp3", False),
    ("EXP002 · Soundtrace 剪", "/Users/renting/Desktop/EXP002_cut_wider.mp3",      True),
]
_positions = [
    (Inches(0.75), Inches(2.2)), (Inches(6.9), Inches(2.2)),
    (Inches(0.75), Inches(4.6)), (Inches(6.9), Inches(4.6)),
]
_cw = Inches(5.7); _ch = Inches(2.15)
for (px, py), (label, path, is_cut) in zip(_positions, audios):
    add_rect(s, px, py, _cw, _ch, fill=C_CARD)
    add_top_border(s, px, py, _cw, Inches(0.05), C_GREEN if is_cut else C_INK3)
    # 标签
    add_text(s, px + Inches(0.3), py + Inches(0.2), _cw - Inches(0.6), Inches(0.5),
             label, size=18, bold=True, color=(C_GREEN if is_cut else C_INK))
    # 副标 · 提示
    add_text(s, px + Inches(0.3), py + Inches(0.75), _cw - Inches(0.6), Inches(0.4),
             ("剪辑后 · 干净的成品" if is_cut else "原始素材 · 未剪"),
             size=13, color=C_INK3)
    # 音频嵌入 · 放中间偏右
    add_audio(s, path, px + Inches(0.4), py + Inches(1.3),
              width=Inches(0.55), height=Inches(0.55))
    # 播放提示
    add_text(s, px + Inches(1.15), py + Inches(1.35), _cw - Inches(1.5), Inches(0.4),
             "▶ 点扬声器图标播放", size=12, color=C_INK3, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.75), Inches(6.85), Inches(11.83), Inches(0.35),
         "两组 A/B 对比 · 无字幕 · 只用耳朵判断", size=12, color=C_INK3)
add_page_num(s, 3)
add_notes(s, """现在放两段音频 · 请大家对比一下。第一组是 EXP019 —— A 是原始 · B 是 Soundtrace 剪出来的。第二组 EXP002 同样。中间没有字幕 · 请大家凭听感判断。听完就大致明白 · 这一整套系统是为了解决什么问题。""")

# ---- P04 · 一键跑通证据 · EP05 首跑 (NEW · 用户 2026-08-21 添加) ----
s = add_slide()
add_sec_chip(s, "§ 1 · 一键跑通证据")
add_h2(s, "一条命令 · 从原始三轨到成品 mp3")
# 命令行框 · 顶部
add_rect(s, Inches(0.75), Inches(1.85), Inches(11.83), Inches(0.75), fill=RGBColor(0x29, 0x25, 0x24))
add_text(s, Inches(1.0), Inches(1.85), Inches(11.5), Inches(0.75),
         "$ python3 run_end_to_end.py --from-raw-wav track_02.wav track_03.wav --episode-id EP05",
         size=13, color=RGBColor(0xFB, 0xBF, 0x24), font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
# 输入 → 输出 大对比
# 左 · 输入
add_rect(s, Inches(0.75), Inches(3.0), Inches(5.4), Inches(3.15), fill=C_CARD)
add_top_border(s, Inches(0.75), Inches(3.0), Inches(5.4), Inches(0.06), C_INK3)
add_text(s, Inches(1.0), Inches(3.2), Inches(5), Inches(0.5), "输入 · 原始",
         size=15, bold=True, color=C_INK3)
add_text(s, Inches(1.0), Inches(3.8), Inches(5), Inches(2.3),
         "3 条对齐 mono wav\n· track_01 (无声 · 已自动跳过)\n· track_02 · 主持人\n· track_03 · 嘉宾\n\n原时长: 26:02 (1562.4 s)",
         size=14, color=C_INK2, line_spacing=1.6)
# 箭头
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.3), Inches(4.35), Inches(0.85), Inches(0.55))
arrow.fill.solid(); arrow.fill.fore_color.rgb = C_ACCENT
arrow.line.fill.background(); arrow.shadow.inherit = False
# 右 · 输出
add_rect(s, Inches(7.25), Inches(3.0), Inches(5.35), Inches(3.15), fill=RGBColor(0xDC, 0xFC, 0xE7))
add_top_border(s, Inches(7.25), Inches(3.0), Inches(5.35), Inches(0.06), C_GREEN)
add_text(s, Inches(7.5), Inches(3.2), Inches(5), Inches(0.5), "输出 · 成品",
         size=15, bold=True, color=C_GREEN)
add_text(s, Inches(7.5), Inches(3.8), Inches(5), Inches(2.3),
         "machine_assisted_draft.mp3\n35 MB · 25:37 · 48 kHz\n\nLLM 主导 · 剪 65 条 = 65 秒\n· filler 38 · repetition 16\n· self_correction 6 · semantic 5",
         size=14, color=C_INK, line_spacing=1.6)
# 底部里程 chip
add_rect(s, Inches(0.75), Inches(6.35), Inches(11.83), Inches(0.65), fill=RGBColor(0xFE, 0xF3, 0xC7))
add_text(s, Inches(0.75), Inches(6.35), Inches(11.83), Inches(0.65),
         "全程无人工干预 · 顺带产出人审判决表 + 相似历史案例 + NISQA 打分建议",
         size=14, bold=True, color=RGBColor(0x78, 0x35, 0x0F),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_page_num(s, 4)
add_notes(s, """这是 EP05 昨晚的实测跑通证据。左边是原料 —— 三条对齐 mono wav · track_01 没声音自动跳过了 · 剩下 track_02 主持人和 track_03 嘉宾。一条命令扔进去。右边是产出 —— 25 分 37 秒的 mp3 · LLM 主导剪掉 65 条 · 涵盖四类语义问题。全程无人工干预 · 顺带产出人审判决表和相似历史案例。这是 "一键输入输出" 的完整证据 · 不是宣传口径 · 是可复现的 run。""")

# ---- P05 · 一键跑通证据 · 聊天记录截图 (NEW · 2026-08-21 v3) ----
s = add_slide()
add_sec_chip(s, "§ 1 · 一键跑通证据")
add_h2(s, "现场跑通 · 聊天记录实录")
# 竖屏截图 · 保比例居中 · 原图 832×1512 (aspect 0.55)
# 页面可用高约 4.8" (从 h2 底部 1.8 到底部 6.6) · 图设置高 4.5" · 宽 4.5*0.55 = 2.48"
_img_h = Inches(4.6)
_img_w = Inches(4.6 * 832 / 1512)  # ~2.53"
_img_x = (SW - _img_w) / 2  # 居中
s.shapes.add_picture(
    "/Users/renting/Desktop/minglue/剪辑项目/交付/最终交付文档/新产出/EP05-PRE-HTML-2026-8-21/p05_evidence_screenshot.jpg",
    _img_x, Inches(1.85), width=_img_w, height=_img_h
) if False else add_picture_safe(s,
    "/Users/renting/Desktop/minglue/剪辑项目/交付/最终交付文档/新产出/EP05-PRE-HTML-2026-8-21/p05_evidence_screenshot.jpg",
    _img_x, Inches(1.85), width=_img_w, height=_img_h)
# 右侧解说
_text_x = _img_x + _img_w + Inches(0.4)
add_text(s, _text_x, Inches(2.0), SW - _text_x - Inches(0.75), Inches(0.5),
         "这就是那次跑通",
         size=20, bold=True, color=C_INK)
add_text(s, _text_x, Inches(2.7), SW - _text_x - Inches(0.75), Inches(4),
         '· 用户丢一句 "按流程跑一遍"\n· pipeline 卡了两次 (Stage 3.5 死锁 / build_calibration_package)\n· 我自己诊断 · 换配置 · 写 driver 绕过\n· 最终 LLM veto 65 KEEP_CUT · 出 25:37 mp3\n\n"不要停 · 我只要成品" —— 用户原话\n全程无需人工介入 · 只交付结果',
         size=13, color=C_INK2, line_spacing=1.7)
add_page_num(s, 5)
add_notes(s, """这一页是那次跑通的原始聊天记录。用户丢一句"按最终交付里的流程跑一遍 · 只要成品"我就上了。中间 pipeline 卡了两次 —— 一次 Stage 3.5 死锁 · 一次 build_calibration_package 老坑 —— 我自己诊断、换配置、写 driver 绕过。最终 LLM veto 判 65 条 KEEP_CUT · 出 25 分 37 秒的成品。整段实录证明这套系统的自动化程度已经能扛住真实运行中的意外 · 而不只是理想路径。""")

# ---- P04 (老) → P06 · 三条难点标题 ----
s = add_slide()
add_sec_chip(s, "§ 2 · 背景")
add_h2(s, '后期 ≠ 只剪 "呃" · 真难点有三条')
pains = [
    ("难点 · 一", "一刀一刀剪\n每刀都要过渡", 'butt splice 会"咔"一下 · 每处都得 crossfade 糊过去 · 过渡长度还挑辅音'),
    ("难点 · 二", "公司音频\n不能上云", "未发布对话属内部数据 · Descript / Adobe Podcast 都得上传 · 我们本地跑"),
    ("难点 · 三", "语义判断\n只有人能做", '"废话还是伏笔？尴尬还是节奏？" · AI 只做手 · 不替代品味'),
]
x = Inches(0.75)
w = Inches(3.9)
for idx, lede, sub in pains:
    add_rect(s, x, Inches(2.2), w, Inches(4.4), fill=C_CARD)
    add_top_border(s, x, Inches(2.2), w, Inches(0.08), C_ACCENT)
    add_text(s, x + Inches(0.3), Inches(2.5), w - Inches(0.6), Inches(0.5), idx,
             size=13, bold=True, color=C_ACCENT)
    add_text(s, x + Inches(0.3), Inches(3.1), w - Inches(0.6), Inches(1.5), lede,
             size=22, bold=True, color=C_INK, line_spacing=1.3)
    add_text(s, x + Inches(0.3), Inches(4.9), w - Inches(0.6), Inches(1.5), sub,
             size=15, color=C_INK3, line_spacing=1.5)
    x += Inches(4.15)
add_page_num(s, 4)
add_notes(s, """一期节目 30 到 60 分钟深度对话 · 主持人加两位嘉宾 · 每个人一支话筒同时录 · 出来是三条对齐的声音轨。大家可能觉得后期不就是把"呃"这种口癖剪掉吗 · 不是。真正的难点有三条 —— 一刀一刀剪都要过渡 · 公司音频不能上云 · 语义判断只有人能做。先说第一条。""")

# ---- P05 · 难点 1 · crossfade ----
s = add_slide()
add_sec_chip(s, "§ 2 · 难点一")
add_h2(s, '直拼会"咔" · 每刀都得 crossfade 糊过去')
# 左卡 · butt splice
add_rect(s, Inches(0.75), Inches(2.2), Inches(5.9), Inches(3.5), fill=C_CARD,
         line=C_ACCENT, line_w=2)
add_text(s, Inches(1.0), Inches(2.4), Inches(5.5), Inches(0.5),
         "butt splice (直拼)", size=18, bold=True, color=C_ACCENT)
# 简易波形示意
line1 = s.shapes.add_connector(1, Inches(1.0), Inches(3.5), Inches(3.5), Inches(3.5))
line1.line.color.rgb = C_INK2; line1.line.width = Pt(2)
line2 = s.shapes.add_connector(1, Inches(3.6), Inches(4.2), Inches(6.2), Inches(4.2))
line2.line.color.rgb = C_INK2; line2.line.width = Pt(2)
break_line = s.shapes.add_connector(1, Inches(3.55), Inches(3.1), Inches(3.55), Inches(4.5))
break_line.line.color.rgb = C_ACCENT; break_line.line.width = Pt(3); break_line.line.dash_style = 5
add_text(s, Inches(1.0), Inches(4.7), Inches(5.5), Inches(1), '断口直接对接 · 十有八九听得出"咔"一下',
         size=14, color=C_INK3)
# 右卡 · crossfade
add_rect(s, Inches(6.9), Inches(2.2), Inches(5.7), Inches(3.5), fill=C_CARD,
         line=C_GREEN, line_w=2)
add_text(s, Inches(7.15), Inches(2.4), Inches(5.3), Inches(0.5),
         "crossfade (交叉淡化)", size=18, bold=True, color=C_GREEN)
line3 = s.shapes.add_connector(1, Inches(7.15), Inches(3.5), Inches(9.5), Inches(3.5))
line3.line.color.rgb = C_INK2; line3.line.width = Pt(2)
line4 = s.shapes.add_connector(1, Inches(9.4), Inches(4.2), Inches(12.3), Inches(4.2))
line4.line.color.rgb = C_GREEN; line4.line.width = Pt(2)
# 重叠区域
overlap = add_rect(s, Inches(9.4), Inches(3.15), Inches(0.5), Inches(1.1), fill=RGBColor(0xD1, 0xFA, 0xE5))
overlap.fill.transparency = 0.4
add_text(s, Inches(7.15), Inches(4.7), Inches(5.3), Inches(1), "几十毫秒曲线 · 断口糊过去 · 听不出剪辑",
         size=14, color=C_INK3)
# 底部说明
add_text(s, Inches(0.75), Inches(6.0), Inches(11.83), Inches(1),
         '过渡长度还得看下一个词的声母 · 太长吃掉辅音("确" → "ue 保")· 太短糊不干净 · 一集要做几十上百次',
         size=15, color=C_INK3, line_spacing=1.5)
add_page_num(s, 5)
add_notes(s, """剪辑者要把问题一个一个找出来 · 每处删除都是一刀 · 每一刀在波形上就是一个断口。断口直接对接叫 butt splice · 中文叫直拼 · 十有八九能听得出"咔"一下 —— 这也是 AI 早期一个典型错误。所以专业剪辑师每一刀都要做一次淡入淡出 · 几十毫秒的曲线把断口糊过去 · 这个手法叫 crossfade。过渡长度还挑辅音 · 太长会把"确"这种 q 声母的词吃成"ue 保"。一集要做几十上百次。""")

# ---- P06 · 难点 2 · 本地 ----
s = add_slide()
add_sec_chip(s, "§ 2 · 难点二")
add_h2(s, "公司音频 · 拔了网线也能跑完")
# 左 · 云工具
add_text(s, Inches(1.0), Inches(2.3), Inches(4.5), Inches(0.5), "云工具",
         size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
y = Inches(3.0)
for tool in ["Descript", "Adobe Podcast", "其它 SaaS"]:
    add_rect(s, Inches(1.5), y, Inches(3.5), Inches(0.65), fill=C_CARD, line=RGBColor(0xFE, 0xCA, 0xCA), line_w=1)
    add_text(s, Inches(1.5), y, Inches(3.5), Inches(0.65), tool,
             size=17, color=C_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.85)
add_text(s, Inches(1.0), Inches(5.7), Inches(4.5), Inches(0.5),
         "都要上传 · 数据出机器", size=12, color=C_INK3, align=PP_ALIGN.CENTER)
# 中箭头
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.8), Inches(3.9), Inches(1.5), Inches(0.6))
arrow.fill.solid(); arrow.fill.fore_color.rgb = C_INK3
arrow.line.fill.background()
arrow.shadow.inherit = False
# 右 · Soundtrace
add_text(s, Inches(7.8), Inches(2.3), Inches(4.5), Inches(0.5), "Soundtrace",
         size=16, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
add_rect(s, Inches(8.0), Inches(3.0), Inches(4), Inches(2.5), fill=RGBColor(0xDC, 0xFC, 0xE7),
         line=C_GREEN, line_w=2)
add_text(s, Inches(8.0), Inches(3.3), Inches(4), Inches(1), "本地跑",
         size=32, bold=True, color=C_INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(8.0), Inches(4.6), Inches(4), Inches(0.7), "全流程离线",
         size=16, color=RGBColor(0x16, 0x65, 0x34), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_page_num(s, 6)
add_notes(s, """未发布的对话内容属内部数据。所有云端在线工具 —— Descript、Adobe Podcast 这些 —— 都要把数据上传。我们的 Soundtrace 全部在自己电脑上完成 · 拔掉网线也能跑完整个流程。这一条不是技术选型问题 · 是合规底线 —— M2 元规则里公司音频永远不出本地。""")

# ---- P07 · 难点 3 · 只有人 ----
s = add_slide()
add_sec_chip(s, "§ 2 · 难点三")
add_h2(s, "语义判断只有人能做 · AI 做手 · 不做品味")
add_left_border(s, Inches(0.85), Inches(2.5), Inches(0.08), Inches(2.5), C_ACCENT)
add_text(s, Inches(1.2), Inches(2.5), Inches(11), Inches(2.5),
         '"这句话是废话还是伏笔？\n这个停顿是尴尬还是节奏？"',
         size=34, bold=True, color=C_INK, line_spacing=1.5)
add_text(s, Inches(1.2), Inches(5.5), Inches(11), Inches(0.5),
         "Mentor 艳馨姐的剪辑作 gold · 系统按人审结果继续学",
         size=14, color=C_INK3)
add_page_num(s, 7)
add_notes(s, """"这句话是废话还是伏笔 · 这个停顿是尴尬还是节奏"—— 这是耳朵和品味的事 · 全用 AI 得不到个性化的结果。我的 mentor 艳馨姐之前就负责音频剪辑 · 我把她的剪辑成果当做 gold · 在"剪哪里"这一层权重最高。系统也会根据用户的每次审核结果继续学。目标是让 AI 做他的手 · 不替代他的品味。""")

# ---- P08 · 6 层俯视 · 移到流程图之后 (2026-08-21 v3 · 匹配用户手改顺序) ----

# ---- P09 · 流程图 1 · 准备层 ----
def build_flow(section_chip, title, stages, page_num, notes, band_color):
    s = add_slide()
    # top band chip · 流程图 X / 4
    add_rect(s, Inches(0.75), Inches(0.55), Inches(1.6), Inches(0.35), fill=band_color)
    add_text(s, Inches(0.75), Inches(0.55), Inches(1.6), Inches(0.35), section_chip,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.5), Inches(0.4), Inches(10), Inches(0.6), title,
             size=26, bold=True, color=C_INK)
    y = Inches(1.4)
    for entry in stages:
        if isinstance(entry, tuple):
            stage_id, body, kind = (entry + (None,))[:3]
            main = kind == "main"
            gate = kind == "gate"
            y = add_flow_stage(s, y, stage_id, body, color=band_color, main=main, gate=gate)
        elif isinstance(entry, str) and entry.startswith("↓"):
            y = add_arrow_down(s, y, entry[1:].strip())
    add_page_num(s, page_num)
    add_notes(s, notes)
    return s

build_flow(
    "流程图 1 / 4", "准备层 · Stage 1–3.5 · 数据准备",
    [
        ("INPUT", "原始三条对齐 mono 声音轨   ·   主持人 + 两位嘉宾 · 各一支话筒"),
        "↓",
        ("S1 · 降噪", "DeepFilterNet v0.5.6 (arm64 CLI)   ·   48 kHz · 尾部 1440 samples 原始回填   ·   --compensate-delay · --atten-lim-db 12"),
        "↓",
        ("S1.5 · 合轨", "3 → 1 主导轨合并   ·   render_prep/speech.mono.wav"),
        "↓",
        ("S2 · ASR", "词级识别 · faster-whisper 1.2.1 · small · int8   ·   每词输出 text · start · end · probability"),
        "↓",
        ("S4 · 分句", "语义分句 · spaCy zh_core_web_sm"),
        "↓",
        ("S5 · 说话人", "说话人分离 · pyannote-audio 4.0.7 · community-1   ·   RTTM · 每段谁在说"),
        "↓",
        ("S6 · 边界精修", "MFA (音素强制对齐) + librosa (onset detect)   ·   剪口精修到 sample 级"),
    ],
    9,
    """原料是三条对齐 mono 声音轨 · 一人一支麦。第一步 DeepFilterNet 做降噪 · 48 千赫 · 尾部 1440 samples 用原始回填 —— 这个补偿参数很关键。第二步合成一条主导轨 · 便于给下游做词级分析。faster-whisper 做词级识别 · 每个词都有 text、start、end、probability 四项。spaCy 做语义分句 · pyannote 4 做说话人分离 · 最后 MFA 加 librosa 把边界精修到 sample 级。词都对好了 · 该判"剪哪些"了。""",
    C_BLUE
)

# ---- P10 · 流程图 2 · PREFERENCE ----
build_flow(
    "流程图 2 / 4", 'PREFERENCE 层 · 判"剪哪些" · LLM 主 · gate 兜',
    [
        ("INPUT", "候选池 · 5 类问题   ·   语气词 / 重复 / 长停顿 / 自纠正 / 串音"),
        "↓",
        ("★ 3.5.5 主", "LLM 语义判决 · Claude Haiku 4.5\n输入每条候选前后 5 秒原文   ·   输出三档: 可以剪 / 别剪保留 / 拿不准交人", "main"),
        '↓ 只放"可以剪"',
        ("3.6 · gate", "结构性硬约束兜底 · 5 门:  说话人身份 · 源轨判定 · 单条时长上限 · 片头片尾 6 秒 · never_cut 硬 override", "gate"),
        "↓",
        ("Stage 5", "EDL 生成 · 只装 verdict==KEEP_CUT   ·   整数 sample · 三轨同步 (M4 元规则)"),
    ],
    10,
    """候选池分五类 —— 语气词、重复、长停顿、自纠正、串音。真正的主决策在 3.5.5 这一步 —— Claude Haiku 4.5 subagent 直接读每条候选前后 5 秒的原文 · 判"可以剪 / 别剪保留 / 拿不准交人"三档 · 只有"可以剪"能进下一步。gate 只做结构性硬约束的兜底 —— 说话人身份、源轨判定、单条时长、片头片尾 6 秒、never_cut · 一共 5 门。另有队列容量和相似度两个 fallback · 但不算硬门。8 门砍到 5 门就是因为语义活 LLM 全接了。最后 Stage 5 生成 EDL · 整数 sample · 三轨同步。""",
    C_PURPLE
)

# ---- P11 · 流程图 3 · PARAMETER ----
build_flow(
    "流程图 3 / 4", 'PARAMETER 层 · 判"怎么剪" · Optuna 独占',
    [
        ("6.7 · 主", "Optuna TPE 贝叶斯优化 · 每候选最多 10 次内收敛\n5 维: crossfade_ms · post_cut_pause_ms · asymmetric_head_pad · boundary_offset_ms · room_tone_pad_ms", "main"),
        "↓",
        ("warm-start", "热启动 · 先跑 2 次已知好起点再自由探索\nA · mentor 59 条中位数    B · YouTube 6 条视频规则   (目前只接 crossfade + room_tone)"),
        "↓",
        ("loss 三项", "signal · 5.0 − discontinuity    +    反作弊 · 漏词 +5.0    +    兜底 · pre_check 挂 +2.0"),
        "↓",
        ("6.9 · 二轮", "人审拒了才触发 · 不用 warm-start · 冷启动重学 10 次 · 补探第一轮没覆盖的区域"),
        "↓",
        ("6.10 · 交集门", '只有 (LLM 说"能剪") ∩ (Optuna 收敛) · 才真正 apply   ·   EP05 首跑实测交集为空 → 跳过 re-render · 门真挡住了', "main"),
    ],
    11,
    """五维参数 · Optuna 用贝叶斯优化去搜 —— crossfade 过渡长度、post_cut_pause 剪后停顿、asymmetric_head_pad 边界扩展、boundary_offset 边界微调、room_tone_pad 环境音补。每候选最多跑 10 次收敛。热启动这一步很关键 —— 先跑 2 次已知好起点再自由探索 · 来源 A 是 mentor 59 条参数的中位数 · 来源 B 是 YouTube 6 条视频规则派生的 · 目前只接了 crossfade 和 room_tone 两维。loss 三项后面单开一页详说。人审拒了会触发 6.9 二轮 · 换新起点冷启动再学 10 次。最后 6.10 交集门 —— 只有 LLM 说"能剪"和 Optuna 收敛的候选才真正 apply。EP05 首跑实测交集为空 · 这个门真挡住了。参数敲定 · 就进 QA。""",
    C_ORANGE
)

# ---- P12 · 流程图 4 · QA ----
build_flow(
    "流程图 4 / 4", "QA + 学习层 · 出成品 · 沉淀经验",
    [
        ("Stage 7 · 混音", "三轨降噪音按 EDL 切好 · 叠成主导轨   ·   -22.2 LUFS 广播标准 · 双遍 loudnorm · 出 mp3"),
        "↓",
        ("Stage 7.1", '整片音质自动打分 · NISQA (德国研究所开源)   ·   达标 → 放行   未达 → 标"建议 mentor 重点听"'),
        "↓",
        ("Stage 7.2", '打包送审 · mp3 + 每条决定理由 + 相似历史案例   ·   markdown 报告 + 待勾选判决表 · 标"机器初稿 · 非最终"'),
        "↓",
        ("沉淀链路", "越用越懂你 · 学习层:  1) mentor 判断存进经验库    2) 立刻重跑候选看翻案    3) 下一期开工自动带进"),
    ],
    12,
    """三轨降噪音轨按剪辑清单切好 · 叠成一条主导轨 · 双遍 loudnorm 压到 -22.2 LUFS 广播行业标准 · 出 mp3。Stage 7.1 用 NISQA 给整片打分 · 这是德国研究所开源的模型 · 达标直接放行 · 不够就标一句"建议 mentor 重点听"。打包时 mp3 加每条决定的理由加相似历史案例 · 做成一份 markdown 报告加一张待勾选判决表 · 明确标"机器初稿 · 非最终"送审。mentor 每次判断都存进经验库 · 立刻重跑一次候选看有没有翻案 · 下一期开工新经验自动带进流程。整条流水线讲完了 · 接下来说控住它的六条元规则。""",
    C_DGREEN
)

# ---- P14 · 6 层俯视 · 移到流程图之后 (2026-08-21 v3 匹配用户手改) ----
s = add_slide()
add_sec_chip(s, "§ 3 · 框架")
add_h2(s, '7 个 skill · 6 层 · 按"边界不能跨"分')
layers = [
    ("L0 · 门卫", "episode-triage", "规则冻结在 plan.json · 不在文档里 —— agent 干着干着会忘文档", C_GREEN),
    ("L1 · 记忆", "feedback-engine", 'session_feedback 单一 SOT · 防"每一集都在重新教机器"', RGBColor(0x08, 0x91, 0xB2)),
    ("L2 · 剪哪些", "PREFERENCE 层 × 3", "candidate-semantic-veto (LLM) 主决策 · gate 只兜结构硬约束", C_PURPLE),
    ("L3 · 怎么剪", "PARAMETER 层 · Optuna", "NISQA discontinuity 作 loss · warm-start 从 mentor + YouTube 派生", C_ORANGE),
    ("L4 · 学经验", "learning-and-experience", "跨 episode 沉淀 · 硬边界:学习层只能写数据源 · 不能改生产代码", RGBColor(0xCA, 0x8A, 0x04)),
    ("L5 · 治理", "governance-and-tool-registry", "未登记的 tool 一律 FAIL · M5 契约先行的执行者", RGBColor(0x6B, 0x72, 0x80)),
]
y = Inches(2.0)
for tag, name, job, color in layers:
    add_rect(s, Inches(0.75), y, Inches(11.83), Inches(0.75), fill=C_CARD)
    add_left_border(s, Inches(0.75), y, Inches(0.06), Inches(0.75), color)
    add_text(s, Inches(1.0), y, Inches(1.9), Inches(0.75), tag,
             size=15, bold=True, color=C_INK2, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.9), y, Inches(2.8), Inches(0.75), name,
             size=13, color=C_INK3, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.75), y, Inches(6.75), Inches(0.75), job,
             size=13, color=C_INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
    y += Inches(0.83)
add_page_num(s, 14)
add_notes(s, """流程走完 · 现在把整套 skill 架构俯视一下。这些 skill 不是按流程顺序分的 · 是按"边界不能越"分的。L0 门卫把关键约束冻结在 plan.json · 不在文档里 · 因为 agent 干着干着会把文档忘了。L1 记忆中枢防"每一集都在重新教机器"。L2 判"剪哪些"归 LLM · L3 判"怎么剪"归 Optuna · 各管各的、不重复不打架 —— 这是 08-20 立的元原则。L4 学经验只能写数据源 · 不能改生产代码。L5 治理让未登记的 tool 一律 FAIL。""")

# ---- P15 (老 P13) · 元规则 ----
s = add_slide()
add_sec_chip(s, "§ 3 · 项目宪法")
add_h2(s, "六条元规则 · 每一条都是踩坑立的")
# M0 通栏
add_rect(s, Inches(0.75), Inches(1.85), Inches(11.83), Inches(1.3), fill=RGBColor(0xFE, 0xFC, 0xE8))
add_top_border(s, Inches(0.75), Inches(1.85), Inches(11.83), Inches(0.06), RGBColor(0xCA, 0x8A, 0x04))
add_text(s, Inches(1.0), Inches(1.95), Inches(11.5), Inches(0.4),
         "M0 · 最高规则 · 开发者模式 · 覆盖 M1-M6",
         size=17, bold=True, color=RGBColor(0xA1, 0x62, 0x07))
add_text(s, Inches(1.0), Inches(2.4), Inches(11.5), Inches(0.4),
         '开发者可临时绕过 M1-M5 · 但 "音频不出本地" 和 "原始素材只读" 永远不能碰',
         size=13, bold=True, color=RGBColor(0x71, 0x3F, 0x12))
add_text(s, Inches(1.0), Inches(2.75), Inches(11.5), Inches(0.4),
         "这次实操让我发现 · 做项目多数时候需要的是修补思维 · 不是搭建思维 · 过多约束会阻碍推进",
         size=12, color=RGBColor(0x71, 0x3F, 0x12))
# M1-M6 · 2x3 grid
m_rules = [
    ("M1", "分层", "生产 / 实验 / 每期 严格隔离", "Champion 只读 · Challenger 隔离 · run-local 只写自己"),
    ("M2", "只读", "原始素材 · 公司音频不出机器", "默认只读 · 不 curl|sh · 不覆盖系统 Python"),
    ("M3", "人签字", "语义删剪必须真人批准", "禁自批准 · 禁超时批准 · 禁伪装成人审"),
    ("M4", "整数采样", "批准区间同步作用三轨", "一个 sample 都不能偏 · 三轨同步整数点"),
    ("M5", "契约先行", "工具必须登记 + adapter", "未登记 → verify.sh 直接 FAIL"),
    ("M6", "报告纪律", "说话分三档 · 已验证/已决定/待验证", '禁用"完成 / 已跑通 / 风险为零"无据词'),
]
cw = Inches(3.8); ch = Inches(1.55)
positions = [(Inches(0.75), Inches(3.35)), (Inches(4.77), Inches(3.35)), (Inches(8.8), Inches(3.35)),
             (Inches(0.75), Inches(5.05)), (Inches(4.77), Inches(5.05)), (Inches(8.8), Inches(5.05))]
for (x, y), (mid, mt, mdt, mds) in zip(positions, m_rules):
    add_rect(s, x, y, cw, ch, fill=C_CARD)
    add_top_border(s, x, y, cw, Inches(0.05), C_ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(0.1), cw - Inches(0.3), Inches(0.35),
             f"{mid} · {mt}", size=15, bold=True, color=C_INK)
    add_text(s, x + Inches(0.2), y + Inches(0.5), cw - Inches(0.3), Inches(0.45), mdt,
             size=12, bold=True, color=C_INK2)
    add_text(s, x + Inches(0.2), y + Inches(1.0), cw - Inches(0.3), Inches(0.5), mds,
             size=10, color=C_INK3, line_spacing=1.4)
add_page_num(s, 13)
add_notes(s, """顶上是六条元规则 —— M0 到 M6 · 我踩坑立的项目宪法。M0 是最高规则 —— 开发者模式可以覆盖 M1 到 M5 · 但音频不出本地和原始素材只读永远不能碰。M0 是我这次实操立的 —— 我发现做项目多数时候需要修补思维 · 不是搭建思维 · 过多约束会阻碍推进。M1 到 M5 依次是分层、只读、人签字、整数采样、契约先行。M6 报告纪律要求说话分三档 —— 已验证事实、已决定方向、待验证假设。每一条都是踩坑立的。下面进入一周实录。""")

# ---- P16 · 踩坑段过场 · token 消耗图 (NEW · 2026-08-21 v3) ----
s = add_slide()
add_sec_chip(s, "§ 4 · 踩坑经验")
add_h2(s, "现在想向大家介绍的是我为期一周的踩坑经历")
# 中部 · token 消耗柱状图 · 原图 688×368 (aspect 1.87)
# 图目标高 3.6" · 宽 3.6 * 688/368 = 6.73"
_ph = Inches(3.6)
_pw = Inches(3.6 * 688 / 368)
_px = (SW - _pw) / 2
add_picture_safe(s,
    "/Users/renting/Desktop/minglue/剪辑项目/交付/最终交付文档/新产出/EP05-PRE-HTML-2026-8-21/p16_token_burn.png",
    _px, Inches(2.4), width=_pw, height=_ph
)
# 图注
add_text(s, Inches(0.75), Inches(6.15), Inches(11.83), Inches(0.4),
         "8 月 12 → 8 月 19 · 每日 output token 消耗 · 接近一周未见实质突破",
         size=14, color=C_INK3, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.75), Inches(6.6), Inches(11.83), Inches(0.4),
         '直到周三下午 · 4 个 Challenger 同时晋升到主流水线 · 局面才转变',
         size=13, color=C_INK2, align=PP_ALIGN.CENTER)
add_page_num(s, 16)
add_notes(s, """这张图是我从 8 月 12 号跑通整个流程到前天终于把最关键的问题解决 · 烧掉的 output token。8/12 六千多 · 8/13 两万六 · 中间几天试来试去没进展 —— 甚至 8/16 一天没动。直到 8/18 三万八、8/19 六万四 · 那两天才是真正跑通的节点。接近一周的时间我没有任何实质性突破。直到周三下午 · 同时有 4 个实验组件从实验区晋升到主流水线 —— pyannote、NISQA、Optuna、案例向量检索 —— 全部默认开启 · 局面才转变。下面这三段我一个一个讲。""")

# ---- P17 (老 P14) · 踩坑 1 ----
s = add_slide()
add_sec_chip(s, "§ 4 · 踩坑 一")
add_h2(s, "参数学出来了 · 但没人用")
# 断链示意
nodes = ["Optuna\n学出参数", "写入\n中间文件", "EDL\n消费", "渲染\n阶段"]
node_colors = [C_INK2, C_INK2, C_ACCENT, C_ACCENT]
node_style = [False, False, True, True]  # dead
nx = Inches(0.85); ny = Inches(2.5); nw = Inches(2.3); nh = Inches(1.2)
gap = Inches(0.6)
for i, (label, col, dead) in enumerate(zip(nodes, node_colors, node_style)):
    x = nx + (nw + gap) * i
    add_rect(s, x, ny, nw, nh, fill=C_CARD, line=col, line_w=2)
    add_text(s, x, ny, nw, nh, label, size=14, bold=True, color=col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
# 箭头 · 前两个 OK · 中间断裂 · 后一个 OK
for i in range(3):
    ax = nx + nw + Emu(0) + (nw + gap) * i - Inches(0.55)
    ay = ny + Inches(0.5)
    if i == 1:
        # broken
        add_text(s, ax, ay - Inches(0.15), Inches(0.6), Inches(0.4), "✗ ✗ ✗",
                 size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    else:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, Inches(0.55), Inches(0.25))
        ar.fill.solid(); ar.fill.fore_color.rgb = C_GRAY
        ar.line.fill.background()
        ar.shadow.inherit = False
add_text(s, Inches(0.85), Inches(4.5), Inches(11.83), Inches(2.5),
         '"接口留着但没人吃"的死代码。听 A/B 不满意 · 让 AI 反复返工 · 逻辑看着都对 —— 实际是"学出来 → 用起来"这条链断了。\n\n修法 · 加一个自动回写环节 · 参数一学出来就强制生效 · 不再走"人审 apply"这一步。',
         size=16, color=C_INK2, line_spacing=1.6)
add_page_num(s, 14)
add_notes(s, """EP04 跑完之后 · 参数自学习学出了一组新参数 · 保存在一个中间文件里 · 但没写回到最终生效的剪辑决策清单 · 渲染阶段用的还是默认值。我听 A/B 对比很不满意 · 看到 AI 给我的流程逻辑都对 · 就让它反复返工。问题的本质是"学出来 → 用起来"这条链断了 —— 一堆接口留着但没人吃的死代码。我后来让 AI 检查每个模块的结果是否被消费才找到 · 加了自动回写 · 参数一学出来就强制生效。""")

# ---- P15 · 踩坑 2 ----
s = add_slide()
add_sec_chip(s, "§ 4 · 踩坑 二")
add_h2(s, "案例检索只看词 · 场景全丢")
add_text(s, Inches(0.85), Inches(2.3), Inches(11.83), Inches(2.0),
         '"呃"这个字 —— 在句首 · 在句中 · 在紧张停顿之后 · 在流利叙述当中 —— 是完全不同的场景 · 但文字上一模一样',
         size=24, bold=True, color=C_INK, line_spacing=1.4)
add_text(s, Inches(0.85), Inches(4.8), Inches(11.83), Inches(2.2),
         '文字匹配的话 · 60 分钟节目里能查到 20 个"呃" · 但它不知道哪 20 个真的可比。\n\n08-19 加了案例向量检索 · 用"听起来像"的方式匹配 —— Whisper encoder 把音频压成向量 · FAISS 找最像的几条。',
         size=16, color=C_INK2, line_spacing=1.6)
add_page_num(s, 15)
add_notes(s, """"呃"这个字 —— 在句首、在句中、在紧张停顿之后、在流利叙述里 —— 是完全不同的场景 · 但文字上一模一样。用文字匹配的话 · 60 分钟节目里能查到 20 个"呃" · 但机器不知道哪 20 个真的可比。08-19 我加了案例向量检索 —— 用"听起来像"的方式匹配 · 而不是词面。""")

# ---- P16 · Mentor 采访 ----
s = add_slide()
add_sec_chip(s, "§ 4 · 转折点")
add_h2(s, "Mentor 采访 · 两个问题让我重想")
# 两个问题
add_rect(s, Inches(0.75), Inches(2.0), Inches(5.9), Inches(2.2), fill=C_CARD)
add_left_border(s, Inches(0.75), Inches(2.0), Inches(0.06), Inches(2.2), C_INK3)
add_text(s, Inches(1.0), Inches(2.15), Inches(5.5), Inches(0.4), "问 一",
         size=12, bold=True, color=C_INK3)
add_text(s, Inches(1.0), Inches(2.55), Inches(5.5), Inches(0.9),
         "碰到 AI 无法解决的问题怎么办？", size=17, bold=True, color=C_INK, line_spacing=1.3)
add_text(s, Inches(1.0), Inches(3.55), Inches(5.5), Inches(0.6),
         "我答：不断调试 · 检查 · 给反馈", size=13, color=C_INK3)
add_rect(s, Inches(6.85), Inches(2.0), Inches(5.65), Inches(2.2), fill=C_CARD)
add_left_border(s, Inches(6.85), Inches(2.0), Inches(0.06), Inches(2.2), C_ACCENT)
add_text(s, Inches(7.1), Inches(2.15), Inches(5.3), Inches(0.4), "问 二",
         size=12, bold=True, color=C_ACCENT)
add_text(s, Inches(7.1), Inches(2.55), Inches(5.3), Inches(0.9),
         "经过这个 Ainol · 你有什么思考？", size=17, bold=True, color=C_INK, line_spacing=1.3)
add_text(s, Inches(7.1), Inches(3.55), Inches(5.3), Inches(0.6),
         "我答：思考问题的底层原理 · 不被表象所困扰", size=13, color=C_INK3)
# 大 quote
add_left_border(s, Inches(0.85), Inches(4.7), Inches(0.08), Inches(2.2), C_ACCENT)
add_text(s, Inches(1.2), Inches(4.7), Inches(11), Inches(2.3),
         "回到工位 · 我想通了：\n我想让 AI 学的是模式 · skill 能做到 ——\n但剪辑的具体时间点、crossfade 值 · 根本不是模式",
         size=20, bold=True, color=C_INK, line_spacing=1.5)
add_page_num(s, 16)
add_notes(s, """周三中午 Mentor 带我们小组去录博客。她说你们要结合具体案例 · 不能只讲空话。我瞬间没底 —— 我压根没跑通核心学习流程 · 连能不能剪干净都保证不了。整个采访我只能说大致思路。之后 Mentor 问了两个问题 —— 碰到 AI 无法解决的问题怎么办 · 以及经过这个 Ainol 你有什么思考。我第二个问题的回答是 —— 思考问题的底层原理 · 而不是被表象所困扰。回到工位我把矛盾想明白了 —— 我想让 AI 学的是模式 · skill 能做到 · 但剪辑的具体时间点和 crossfade 值根本不是模式。参数和候选一定要分离。""")

# ---- P17 · Optuna loss 三步 ----
s = add_slide()
add_sec_chip(s, "§ 4 · 突围 一 · PARAMETER")
add_h2(s, "Optuna 引入 · loss 是三步推出来的")
add_text(s, Inches(0.85), Inches(1.85), Inches(11.83), Inches(0.5),
         '起点 · 以 NISQA discontinuity 作 benchmark · 5 维打分里的"不连续性" · 满分 5.0 · 越高越干净',
         size=14, color=C_INK3)
steps = [
    ("STEP 1", "镜像", "Optuna 默认最小化目标 · 把 discontinuity 反过来", "loss = 5.0 − discontinuity"),
    ("STEP 2", "反作弊", '发现走捷径 —— 不剪就没剪口 · 天然满分。加惩罚让"必须真剪"', "+ 5.0 × 每漏一个目标词"),
    ("STEP 3", "兜底", "基础检查(静音位置)挂了再加一个较轻的惩罚 · 挡明显破坏", "+ 2.0 (若 pre_check FAIL)"),
]
sx = Inches(0.85); sy = Inches(3.0); sw = Inches(3.95); sh = Inches(3.5); sgap = Inches(0.15)
for i, (n, l, d, f) in enumerate(steps):
    x = sx + (sw + sgap) * i
    add_rect(s, x, sy, sw, sh, fill=C_CARD)
    add_top_border(s, x, sy, sw, Inches(0.06), C_ORANGE)
    add_text(s, x + Inches(0.25), sy + Inches(0.2), sw, Inches(0.3), n,
             size=11, bold=True, color=C_ORANGE)
    add_text(s, x + Inches(0.25), sy + Inches(0.55), sw, Inches(0.5), l,
             size=20, bold=True, color=C_INK)
    add_text(s, x + Inches(0.25), sy + Inches(1.2), sw - Inches(0.4), Inches(1.4), d,
             size=13, color=C_INK3, line_spacing=1.5)
    add_text(s, x + Inches(0.25), sy + Inches(2.7), sw - Inches(0.4), Inches(0.6), f,
             size=14, bold=True, color=C_ORANGE, font=FONT_MONO)
add_page_num(s, 17)
add_notes(s, """之前有个我说不清但听着别扭的现象 —— 规则再多也只能挡"已知的错" · 挡不了这种模糊听感。所以我需要一个自动打分的 signal · 让机器替我判"干不干净"。起点是 NISQA discontinuity —— NISQA 五维打分里的"不连续性"这一维 · 满分 5 · 越高越干净。这就是核心 signal。三步推 loss —— 第一步镜像 · Optuna 默认最小化目标 · 把 discontinuity 反过来写成 5 减 discontinuity。第二步反作弊 —— 发现只用这个 Optuna 会走捷径 · 通过"根本不剪"拿满分 · 所以加第一个 penalty 让它明白必须真剪掉才算成功。第三步兜底 —— 若基础检查挂了再加 2.0 惩罚 · 比漏一个词轻、比正常波动重。loss 就三项 · 定了。""")

# ---- P18 · warm-start ----
s = add_slide()
add_sec_chip(s, "§ 4 · 突围 二 · warm-start")
add_h2(s, "10 个案例 7 个卡在局部最小 · 热启动破局")
add_text(s, Inches(0.85), Inches(1.9), Inches(11.83), Inches(1.2),
         '想到之前打数学建模美赛也碰过局部最小的问题 —— 让 AI 把"相似情况下 mentor 的参数"和"YouTube 上的经验"作为搜索起点 · 先跑 2 次已知好点 · 再自由探索。',
         size=15, color=C_INK2, line_spacing=1.6)
# 两卡
add_rect(s, Inches(0.85), Inches(3.5), Inches(5.85), Inches(2.5), fill=C_CARD)
add_left_border(s, Inches(0.85), Inches(3.5), Inches(0.08), Inches(2.5), C_ORANGE)
add_text(s, Inches(1.15), Inches(3.7), Inches(5.4), Inches(0.5), "来源 A · mentor 手工剪",
         size=17, bold=True, color=C_INK)
add_text(s, Inches(1.15), Inches(4.3), Inches(5.4), Inches(1.5),
         "59 条参数的中位数 · 从 EP03 + EP04 反向提取\ncrossfade / boundary_offset 每类 kind 独立",
         size=14, color=C_INK3, line_spacing=1.6)
add_rect(s, Inches(6.9), Inches(3.5), Inches(5.6), Inches(2.5), fill=C_CARD)
add_left_border(s, Inches(6.9), Inches(3.5), Inches(0.08), Inches(2.5), C_ORANGE)
add_text(s, Inches(7.2), Inches(3.7), Inches(5.2), Inches(0.5), "来源 B · YouTube 6 条规则",
         size=17, bold=True, color=C_INK)
add_text(s, Inches(7.2), Inches(4.3), Inches(5.2), Inches(1.5),
         "专业剪辑视频里派生的经验\n目前只接了 crossfade + room_tone · 其它 3 维仍从零学",
         size=14, color=C_INK3, line_spacing=1.6)
add_text(s, Inches(0.85), Inches(6.3), Inches(11.83), Inches(0.5),
         "收敛率大幅提升 · 我很满意", size=18, bold=True, color=C_GREEN)
add_page_num(s, 18)
add_notes(s, """loss 定了 · 扔给 AI 一跑 —— 10 个案例 7 个不收敛。我想到之前打数学建模美赛也碰过这个问题 —— 卡在局部最小。突然想到可以把相似情况下 Mentor 手工剪的参数、YouTube 上高手的经验作为搜索起点 · 做一个热启动 —— 先跑 2 次已知好点 · 再让 Optuna 自由探索。这也回答了我之前的疑问 —— 有了算法 · Mentor gold 和 YouTube 经验依然有意义 · 它们不做判决 · 但给搜索一个好的入口。收敛率大幅提升 · 我很满意。""")

# ---- P19 · 候选层转向 ----
s = add_slide()
add_sec_chip(s, "§ 4 · 突围 三 · PREFERENCE")
add_h2(s, "候选层从规则派转向 LLM 主决策")
add_text(s, Inches(0.85), Inches(1.9), Inches(11.83), Inches(0.7),
         '失忆症 —— EP04 说"一些"不许剪 · EP05 又剪了一次 · 每期都在重新教',
         size=18, bold=True, color=C_INK)
# 左 · 之前
add_rect(s, Inches(0.85), Inches(3.0), Inches(5.35), Inches(3.4), fill=C_CARD)
add_top_border(s, Inches(0.85), Inches(3.0), Inches(5.35), Inches(0.06), C_ACCENT)
add_text(s, Inches(1.1), Inches(3.2), Inches(5), Inches(0.5), "之前 · 规则派",
         size=17, bold=True, color=C_ACCENT)
for i, line in enumerate(["· 5 类候选族", "· 8 道 gate", "· 整齐 · 但离用户实际差太远", "· 规则永远追不上语义"]):
    add_text(s, Inches(1.1), Inches(3.8 + i * 0.5), Inches(5), Inches(0.45), line,
             size=14, color=C_INK2)
# 中箭头
mid_arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.35), Inches(4.4), Inches(0.65), Inches(0.5))
mid_arrow.fill.solid(); mid_arrow.fill.fore_color.rgb = C_INK3
mid_arrow.line.fill.background()
mid_arrow.shadow.inherit = False
# 右 · 现在
add_rect(s, Inches(7.15), Inches(3.0), Inches(5.35), Inches(3.4), fill=C_CARD)
add_top_border(s, Inches(7.15), Inches(3.0), Inches(5.35), Inches(0.06), C_GREEN)
add_text(s, Inches(7.4), Inches(3.2), Inches(5), Inches(0.5), "现在 · LLM 主 · gate 兜",
         size=17, bold=True, color=C_GREEN)
for i, line in enumerate([
    "· Claude Haiku 4.5 subagent 读原文",
    "· 判 KEEP_CUT / REJECT_KEEP / NEEDS_REVIEW",
    "· gate 从 8 门砍到 5 门 · 只留结构红线",
    "· never_cut 是硬 override · Haiku 也推不翻",
]):
    add_text(s, Inches(7.4), Inches(3.8 + i * 0.5), Inches(5), Inches(0.45), line,
             size=13, color=C_INK2)
# 底部数据
add_rect(s, Inches(0.85), Inches(6.6), Inches(11.65), Inches(0.5), fill=C_CARD)
add_left_border(s, Inches(0.85), Inches(6.6), Inches(0.08), Inches(0.5), C_GREEN)
add_text(s, Inches(1.15), Inches(6.6), Inches(11.3), Inches(0.5),
         "EP05 前五分钟 ·  37 候选 · 33 自动剪 · 4 交人审   ·   自动率 89.2%",
         size=15, bold=True, color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
add_page_num(s, 19)
add_notes(s, """参数收得差不多 · 我转头把候选也收网 —— 就是判"这一段该不该剪"。早期我是规则派 · 5 类候选加 8 道 gate · 觉得很整齐。结果 8-17 用户明确"一些"不许剪 · 下一期系统又剪了一次 —— 我给这毛病起了个名字叫失忆症。我接了 Claude Haiku 4.5 的 subagent 让它直接读原文判 KEEP_CUT / REJECT_KEEP / NEEDS_REVIEW。第一次跑通看到"一些"被标 REJECT_KEEP · 我才觉得机器开始有点像人。8 门 gate 一路砍到 5 门 · 只留结构红线。EP05 前五分钟 37 候选 33 自动剪 4 交人审 · 自动率 89.2%。""")

# ---- P20 · learning-pattern ----
s = add_slide()
add_sec_chip(s, "§ 4 · 新起 skill")
add_h2(s, "learning-pattern-from-case-v1 · 让 LLM 读自然语言")
add_text(s, Inches(0.85), Inches(1.9), Inches(11.83), Inches(0.9),
         "经验卡按 reason_key 分桶 · 每桶按 accept/reject 分 · 挑 3-5 个代表 case · 把 deleted_text 和 evidence 一起喂 LLM · 让它归纳共同语义模式",
         size=15, color=C_INK2, line_spacing=1.5)
# pattern md preview 深底
add_rect(s, Inches(0.85), Inches(3.15), Inches(11.83), Inches(3.15), fill=RGBColor(0x29, 0x25, 0x24))
pattern_lines = [
    ("## reason_key: filler_hesitation", RGBColor(0xFB, 0xBF, 0x24), True),
    ("", None, False),
    ("### decision: reject (mentor 剪掉)", RGBColor(0x93, 0xC5, 0xFD), True),
    ("- 模式: 短促无信息 filler · 单独出现在句间", RGBColor(0xE7, 0xE5, 0xE4), False),
    ('- 例子: [EP03-C007] "然后呃我们..."', RGBColor(0xE7, 0xE5, 0xE4), False),
    ('        [EP03-C014] "所以嗯就是..."', RGBColor(0xE7, 0xE5, 0xE4), False),
    ("", None, False),
    ("### decision: accept (mentor 保留)", RGBColor(0x93, 0xC5, 0xFD), True),
    ("- 模式: 情绪 filler 或语气强调", RGBColor(0xE7, 0xE5, 0xE4), False),
    ('- 例子: [EP03-C023] "呃 这个真的很难说"', RGBColor(0xE7, 0xE5, 0xE4), False),
]
tb = s.shapes.add_textbox(Inches(1.1), Inches(3.3), Inches(11.3), Inches(2.9))
tb.text_frame.word_wrap = True
tb.text_frame.margin_top = Emu(0); tb.text_frame.margin_bottom = Emu(0)
for i, (line, col, bold) in enumerate(pattern_lines):
    p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
    p.line_spacing = 1.35
    r = p.add_run()
    r.text = line or " "
    r.font.name = FONT_MONO
    r.font.size = Pt(12)
    r.font.bold = bold
    if col: r.font.color.rgb = col
add_text(s, Inches(0.85), Inches(6.4), Inches(11.83), Inches(0.9),
         "严格无数字 · 只出自然语言 pattern · 每条挂 case_id 能追溯 · mentor 一眼核对\n硬 checklist 8 条 · 只学 PREFERENCE (剪哪些) · 绝不碰 PARAMETER (怎么剪)",
         size=13, color=C_INK3, line_spacing=1.5)
add_page_num(s, 20)
add_notes(s, """我新起一个 learning-pattern-from-case-v1 skill —— 把 Mentor 手动剪辑的部分和之前的审核记录做成经验卡。按 reason_key 分桶 —— filler、repetition、long_pause、self_correction —— 每桶再按 accept / reject 分 · 每桶挑 3-5 个代表 case · 把 deleted_text 和 evidence_text 一起喂 LLM · 让它归纳共同语义模式 —— 输出一份 pattern md。严格无数字 · 只出自然语言 pattern · 每条挂 case_id 能追溯。LLM 读自然语言比读统计好 —— 这就是这次重写的根。只学 PREFERENCE · 绝不碰 PARAMETER · 硬 checklist 8 条挡在那。""")

# ---- P21 · 未来 4 方向 ----
s = add_slide()
add_sec_chip(s, "§ 5 · 未来")
add_h2(s, "基础技术层收敛 · 重心转向偏好积累")
futures = [
    ("方向 · 一", "案例记忆持续积累", "65 条历史案例 + 33 条人审 + 42 条 YouTube = 140 项历史知识在被消费。gate 第 5 关和 6.8 步向量检索两条通路同时消费。期数越多 · 记忆越贴合 —— 这套是章鱼 AI 独有的"),
    ("方向 · 二", "learning-pattern-from-case-v1 激活", '开出第三条通路 · 把历史案例蒸馏成 LLM 可读的自然语言模式 · 让主流水线判决直接引用"上一期 mentor 是怎么处理类似情况" —— 而不只是相似度打分'),
    ("方向 · 三", "向量数据库 · 已下载在试验区", "Whisper encoder 压成 512 维向量 · FAISS 找最像 · 只出相似度分数 · 不下判决。目前没投入是因为案例量还不够验证 —— 等更多期数积累后启用"),
    ("方向 · 四", "Optuna 搜索空间扩展", '现在只对"剪的干不干净"做参数级迭代 · 未来把响度阈值、闪避深度、音乐过渡曲线也纳进来 · 让参数学习覆盖整条流水线'),
]
fx = Inches(0.85); fy0 = Inches(2.0); fw = Inches(5.85); fh = Inches(2.5); fgap_x = Inches(0.15); fgap_y = Inches(0.15)
positions = [(fx, fy0), (fx + fw + fgap_x, fy0), (fx, fy0 + fh + fgap_y), (fx + fw + fgap_x, fy0 + fh + fgap_y)]
for (x, y), (fn, fl, fd) in zip(positions, futures):
    add_rect(s, x, y, fw, fh, fill=C_CARD)
    add_top_border(s, x, y, fw, Inches(0.05), C_PURPLE)
    add_text(s, x + Inches(0.25), y + Inches(0.15), fw - Inches(0.4), Inches(0.35),
             fn, size=11, bold=True, color=C_PURPLE)
    add_text(s, x + Inches(0.25), y + Inches(0.5), fw - Inches(0.4), Inches(0.5),
             fl, size=17, bold=True, color=C_INK)
    add_text(s, x + Inches(0.25), y + Inches(1.1), fw - Inches(0.4), fh - Inches(1.2),
             fd, size=12, color=C_INK3, line_spacing=1.55)
add_page_num(s, 21)
add_notes(s, """MFA 对齐、剪口检查五项、Optuna 自学习、案例向量检索这些都落地之后 · 基础技术层基本收敛。重心从"技术攻关"进入"偏好积累"。四个方向 —— 一是案例记忆持续积累 · 已经有 65 条历史加 33 条人审加 42 条 YouTube 共 140 项在被消费。二是 learning-pattern-from-case-v1 激活后 · LLM 判决能直接引用"上一期 Mentor 是怎么处理类似情况"的原文。三是向量数据库 · 已下载在试验区 · 等案例量够就上。四是 Optuna 搜索空间扩到响度、闪避、音乐过渡。下面用一句话收尾。""")

# ---- P22 · 收尾 ----
s = add_slide(bg=C_BG_DARK)
add_brand_line(s, Inches(6), Inches(2.0), width=Inches(1.5), height=Inches(0.06))
add_text(s, Inches(0.75), Inches(2.5), Inches(11.83), Inches(3.5),
         '我们没有把播客后期变成一个 "魔法按钮"。\n\n我们做的是 —— 让每一次剪、每一次批准、每一次学习\n都留下痕迹、可以回放、可以被下一期节目用上。\n\n人还在决策的中心 · 机器在旁边帮忙',
         size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, line_spacing=1.5)
add_text(s, Inches(0.75), Inches(6.3), Inches(11.83), Inches(0.6),
         "谢 谢 大 家", size=28, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1),
         align=PP_ALIGN.CENTER)
add_page_num(s, 22, dark=True)
add_notes(s, """我们没有把播客后期变成一个"魔法按钮"。我们做的是 —— 让每一次剪、每一次批准、每一次学习都留下痕迹、可以回放、可以被下一期节目用上。人还在决策的中心 · 机器在旁边帮忙 · 而不是反过来。每一份代码、每一个决策、每一个 skill 的边界都写在纸上 —— 是为了让下一个接手者 · 能够在我离开之后把这套东西继续做下去。谢谢大家。""")

# ==================== 保存 ====================
finalize_page_nums()
prs.save(str(OUT))
print(f"✅ 生成完成 · {OUT}")
print(f"   slides = {len(prs.slides)}")
